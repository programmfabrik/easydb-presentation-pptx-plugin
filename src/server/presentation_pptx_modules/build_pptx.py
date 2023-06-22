# encoding: utf-8

import collections
import collections.abc

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from PIL import Image

from . import pptx_util

import hashlib
import os


def __pixels_to_emu(px, dpi=72):
    return float(px * (914400 / dpi))


def __get_standard_format():
    return {
        '1': {
            'size': 26,
            'bold': True
        },
        '2': {
            'size': 19,
            'bold': False
        },
        '3': {
            'size': 16,
            'bold': False
        }
    }


def __insert_info(text_placeholder, shapes, standard_info, show_standard, standard_format):
    if not isinstance(standard_info, dict):
        return

    text_box = shapes.add_textbox(text_placeholder.left, text_placeholder.top, text_placeholder.width, text_placeholder.height)
    text_box.text_frame.word_wrap = True

    first_standard_value = True
    for s in show_standard:
        if not s in standard_format:
            continue
        if not s in standard_info:
            continue

        _s = standard_info[s].strip()
        if len(_s) < 1:
            continue

        if first_standard_value:
            first_standard_value = False
            p = text_box.text_frame.paragraphs[0]
        else:
            p = text_box.text_frame.add_paragraph()

        p.text = _s
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Helvetica'
        p.font.size = Pt(standard_format[s]['size'])
        p.font.bold = standard_format[s]['bold']

    # remove the original placeholder since it is not needed
    text_placeholder._element.getparent().remove(text_placeholder._element)


def __insert_text(picture_placeholder, placeholder_to_remove, shapes, text):
    # remove the obsolete text placeholder under the picture
    placeholder_to_remove._element.getparent().remove(placeholder_to_remove._element)

    if len(text) < 1:
        return

    lines = []
    first_line = True
    for s in text.split('\n'):
        lines.append(s.strip())

    if len(lines) < 1:
        return

    fontsize = 26
    th = __pixels_to_emu(fontsize * len(lines))
    text_box = shapes.add_textbox(
        left=picture_placeholder.left,
        top=picture_placeholder.top + ((picture_placeholder.height - th) / 2),
        width=picture_placeholder.width,
        height=th
    )
    text_box.text_frame.word_wrap = True

    for s in lines:
        if first_line:
            first_line = False
            p = text_box.text_frame.paragraphs[0]
        else:
            p = text_box.text_frame.add_paragraph()

        p.text = s
        p.alignment = PP_ALIGN.LEFT
        p.font.name = 'Helvetica'
        p.font.size = Pt(fontsize)
        p.font.bold = True

    # remove the original placeholder since it is not needed
    picture_placeholder._element.getparent().remove(picture_placeholder._element)


def __insert_picture(pack_dir, exp_files, picture_placeholder, shapes, eas_id, asset_url, placeholder_image, placeholder_info):

    if eas_id is None and asset_url is None:
        if placeholder_info is None:
            return

    filename = None
    print_placeholder_info = False

    if asset_url is not None:
        try:
            # download the image file, save it in the export asset folder
            m = hashlib.sha1(asset_url)
            filename = '{0}/{1}'.format(pack_dir, str(m.hexdigest()))

            url_parts = asset_url.split('/')
            if len(url_parts) > 1:
                filename += '.{0}'.format(url_parts[-1])

            pptx_util.download_export_file(asset_url, filename)
        except Exception as e:
            print('could not download connector image: {0}'.format(str(e)))
            filename = placeholder_image
            print_placeholder_info = True
    else:
        for _file in exp_files:
            if not 'eas_id' in _file:
                continue
            if _file['eas_id'] != eas_id:
                continue
            filename = os.path.abspath('{0}/{1}'.format(pack_dir, _file['path']))
            break

        if filename is None:
            filename = placeholder_image
            print_placeholder_info = True

    try:
        img = Image.open(filename)
    except Exception as e:
        print('could not load exported image {0}: {1}'.format(filename, str(e)))

    try:
        # get placeholder size in emus
        pw_emu = float(picture_placeholder.width)
        ph_emu = float(picture_placeholder.height)

        iw, ih = img.size

        # convert image size from pixels to emus
        iw_emu = __pixels_to_emu(iw)
        ih_emu = __pixels_to_emu(ih)

        h_ratio = iw_emu / pw_emu
        w_ratio = ih_emu / ph_emu

        # scale down to fit the longer image side into the shorter placeholder side
        # make sure to not stretch smaller images
        if h_ratio >= w_ratio:
            new_h = min(ih_emu, int(ih_emu / h_ratio))
            new_w = min(iw_emu, int(iw_emu / h_ratio))
        else:
            new_h = min(ih_emu, int(ih_emu / w_ratio))
            new_w = min(iw_emu, int(iw_emu / w_ratio))

        new_x = picture_placeholder.left + (pw_emu - new_w) / 2
        new_y = picture_placeholder.top + (ph_emu - new_h) / 2

        shapes.add_picture(filename, new_x, new_y, height=new_h)

        if print_placeholder_info:
            # use the info string about the asset to display something informative along the placeholder
            fontsize = 15
            text_box = shapes.add_textbox(
                left=new_x,
                top=new_y + new_h - __pixels_to_emu(2 * (fontsize + 5)),
                width=new_w,
                height=__pixels_to_emu(fontsize + 5)
            )
            text_box.text_frame.word_wrap = True
            p = text_box.text_frame.paragraphs[0]
            p.text = placeholder_info
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Helvetica'
            p.font.size = Pt(fontsize)

        # remove the original placeholder since it is not needed
        picture_placeholder._element.getparent().remove(picture_placeholder._element)
    except Exception as e:
        picture_placeholder.insert_picture(placeholder_image)


def produce_files(produce_opts, pack_dir, export_files, pptx_filename):

    standard_format = __get_standard_format()
    show_standard = pptx_util.get_json_value(produce_opts, 'presentation.settings.show_standard')
    if isinstance(show_standard, str):
        show_standard = list(map(lambda s: s.strip(), show_standard.split()))
    else:
        show_standard = []

    template = pptx_util.get_json_value(produce_opts, 'pptx_form.template', True)

    cur_dir = os.path.abspath(os.path.dirname(__file__))
    prs = Presentation(os.path.join(cur_dir, '..', '..', 'templates', pptx_util.get_json_value(template, 'name', True)))
    slide_layouts = {}
    for slide in pptx_util.get_json_value(template, 'slides', True):
        slide_layouts[slide['type']] = {
            'layout': prs.slide_layouts[slide['slide_idx']],
            'info': slide,
        }
    placeholder_image = os.path.join(cur_dir, '..', '..', 'placeholders', pptx_util.get_json_value(template, 'placeholder', True))

    slide_id = -1
    for slide in pptx_util.get_json_value(produce_opts, 'presentation.slides', True):
        slide_id += 1

        stype = slide['type']

        sl = slide_layouts[stype]
        sl_info = sl['info']

        ppt_slide = prs.slides.add_slide(sl['layout'])

        title_key = pptx_util.get_json_value(sl_info, 'title')
        subtitle_key = pptx_util.get_json_value(sl_info, 'subtitle')

        data_title = pptx_util.get_json_value(slide, 'data.title')
        data_info = pptx_util.get_json_value(slide, 'data.info')

        if stype == 'start':
            if not 'data' in slide:
                continue

            if title_key is not None and data_title is not None:
                ppt_slide.placeholders[title_key].text = data_title
            if subtitle_key is not None and data_info is not None:
                ppt_slide.placeholders[subtitle_key].text = data_info

        elif stype == 'bullets':
            if not 'data' in slide:
                continue

            if title_key is not None and data_title is not None:
                ppt_slide.placeholders[title_key].text = data_title

            text_frame = ppt_slide.placeholders[sl_info['bullets']].text_frame
            text_frame.clear()  # remove any existing paragraphs, leaving one empty one

            rows = slide['data']['info'].split('\n')

            p = text_frame.paragraphs[0]
            p.text = rows[0]

            for row in rows[1:]:
                p = text_frame.add_paragraph()
                p.text = row

        elif stype == 'one':
            if not 'center' in slide:
                continue

            if not 'global_object_id' in slide['center']:
                continue

            __insert_picture(
                pack_dir,
                export_files,
                ppt_slide.placeholders[sl_info['picture']],
                ppt_slide.shapes,
                pptx_util.get_json_value(slide, 'center.version_id'),
                pptx_util.get_json_value(slide, 'center.asset_url'),
                placeholder_image,
                pptx_util.get_json_value(slide, 'center.placeholder_info'),
            )

            if 'text' in sl_info:
                __insert_info(
                    ppt_slide.placeholders[sl_info['text']],
                    ppt_slide.shapes,
                    pptx_util.get_json_value(slide, 'center.standard_info'),
                    show_standard,
                    standard_format,
                )

        elif stype == 'duo':

            if not 'left' in slide and not 'right' in slide:
                continue

            if 'left' in slide:
                if 'global_object_id' in slide['left'] and 'picture_left' in sl_info:
                    __insert_picture(
                        pack_dir,
                        export_files,
                        ppt_slide.placeholders[sl_info['picture_left']],
                        ppt_slide.shapes,
                        pptx_util.get_json_value(slide, 'left.version_id'),
                        pptx_util.get_json_value(slide, 'left.asset_url'),
                        placeholder_image,
                        pptx_util.get_json_value(slide, 'left.placeholder_info'),
                    )

            if 'right' in slide:
                if 'global_object_id' in slide['right'] and 'picture_right' in sl_info:
                    __insert_picture(
                        pack_dir,
                        export_files,
                        ppt_slide.placeholders[sl_info['picture_right']],
                        ppt_slide.shapes,
                        pptx_util.get_json_value(slide, 'right.version_id'),
                        pptx_util.get_json_value(slide, 'right.asset_url'),
                        placeholder_image,
                        pptx_util.get_json_value(slide, 'right.placeholder_info'),
                    )

            if 'left' in slide:
                if 'global_object_id' in slide['left'] and 'text_left' in sl_info:
                    __insert_info(
                        ppt_slide.placeholders[sl_info['text_left']],
                        ppt_slide.shapes,
                        pptx_util.get_json_value(slide, 'left.standard_info'),
                        show_standard,
                        standard_format,
                    )

            if 'right' in slide:
                if 'global_object_id' in slide['right'] and 'text_right' in sl_info:
                    __insert_info(
                        ppt_slide.placeholders[sl_info['text_right']],
                        ppt_slide.shapes,
                        pptx_util.get_json_value(slide, 'right.standard_info'),
                        show_standard,
                        standard_format,
                    )

        elif stype == 'imageText':

            if not 'left' in slide and not 'data' in slide:
                continue

            if 'left' in slide:
                if 'global_object_id' in slide['left'] and 'picture_left' in sl_info:
                    __insert_picture(
                        pack_dir,
                        export_files,
                        ppt_slide.placeholders[sl_info['picture_left']],
                        ppt_slide.shapes,
                        pptx_util.get_json_value(slide, 'left.version_id'),
                        pptx_util.get_json_value(slide, 'left.asset_url'),
                        placeholder_image,
                        pptx_util.get_json_value(slide, 'left.placeholder_info'),
                    )

            if 'left' in slide:
                if 'global_object_id' in slide['left'] and 'text_left' in sl_info:
                    __insert_info(
                        ppt_slide.placeholders[sl_info['text_left']],
                        ppt_slide.shapes,
                        pptx_util.get_json_value(slide, 'left.standard_info'),
                        show_standard,
                        standard_format,
                    )

            text = pptx_util.get_json_value(slide, 'data.text')
            if isinstance(text, str) and 'text_right' in sl_info:
                __insert_text(
                    ppt_slide.placeholders[sl_info['text_right']],
                    ppt_slide.placeholders[sl_info['text_box_to_remove']],
                    ppt_slide.shapes,
                    text)

    pptx_util.create_missing_dirs(pptx_filename)
    prs.save(pptx_filename)
