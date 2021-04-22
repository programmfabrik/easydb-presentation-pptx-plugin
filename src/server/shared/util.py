# encoding: utf-8

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from PIL import Image

import hashlib
import urllib


def get_json_value(js, path, expected=False):
    current = js
    path_parts = path.split('.')
    for path_part in path_parts:
        if not isinstance(current, dict) or path_part not in current:
            if expected:
                raise Exception('expected: {0}'.format(path_part))
            else:
                return None
        current = current[path_part]
    return current


def pixels_to_emu(px, dpi=72):
    return float(px * (914400 / dpi))


def new_presentation(template_path):
    return Presentation(template_path)


LOG_DEBUG = 1
LOG_INFO = 2
LOG_WARN = 3
LOG_ERROR = 4


def easydb5log(logger, msg, level=LOG_INFO):
    if logger is None:
        return

    if level == LOG_DEBUG:
        logger.debug(msg)
    elif level == LOG_INFO:
        logger.info(msg)
    elif level == LOG_WARN:
        logger.warn(msg)
    elif level == LOG_ERROR:
        logger.error(msg)


def get_standard_format():
    return {
        '1': {
            'size': 26,
            'bold': True
        },
        '2': {
            'size': 19,
            'bold': False
        },
        '3': {
            'size': 16,
            'bold': False
        }
    }


def parse_target_filename(produce_opts):
    return '%s.pptx' % produce_opts['presentation']['filename']


def parse_show_standard(produce_opts):
    if not 'show_standard' in produce_opts['presentation']['settings']:
        return []
    return [s.strip() for s in produce_opts['presentation']['settings']['show_standard'].split(' ')]


def parse_slide_layouts(prs, produce_opts, show_standard):
    slide_layouts = {}

    for slide in produce_opts['pptx_form']['template']['slides']:
        if len(show_standard) > 0:
            if 'show_info' in slide and slide['show_info'] == True:
                slide_layouts[slide['type']] = {
                    'layout': prs.slide_layouts[slide['slide_idx']],
                    'info': slide,
                }
        else:
            if 'show_info' not in slide or slide['show_info'] == False:
                slide_layouts[slide['type']] = {
                    'layout': prs.slide_layouts[slide['slide_idx']],
                    'info': slide,
                }

    return slide_layouts


def insert_info(placeholder, shapes, gid, data_by_gid, show_standard, standard_format, picture_bottom_line):
    if gid not in data_by_gid:
        return

    top = placeholder.top

    if picture_bottom_line is not None:
        top = min([picture_bottom_line + pixels_to_emu(10), top])

    text_box = shapes.add_textbox(
        placeholder.left, top, placeholder.width, placeholder.height)
    text_box.text_frame.word_wrap = True

    first_standard_value = True
    for s in show_standard:
        if not s in standard_format:
            continue

        if s in data_by_gid[gid]['standard_info']:
            _s = data_by_gid[gid]['standard_info'][s].strip()
            if len(_s) < 1:
                continue

            if first_standard_value:
                first_standard_value = False
                p = text_box.text_frame.paragraphs[0]
            else:
                p = text_box.text_frame.add_paragraph()

            p.text = _s
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Helvetica'
            p.font.size = Pt(standard_format[s]['size'])
            p.font.bold = standard_format[s]['bold']

    # remove the original placeholder since it is not needed
    placeholder._element.getparent().remove(placeholder._element)


def insert_picture(exp_files_path, exp_files, placeholder, shapes, eas_id, asset_url, logger=None):

    picture_bottom_line = None

    if eas_id is None and asset_url is None:
        easydb5log(logger,
                   'no asset id or asset url is given for insert_picture',
                   LOG_WARN)
        return

    filename = None
    use_connector_url = False

    if asset_url is not None:
        try:
            # download the image file, save it in the export asset folder
            m = hashlib.sha1(asset_url)
            filename = '%s/%s' % (exp_files_path, str(m.hexdigest()))

            url_parts = asset_url.split('/')
            if len(url_parts) > 1:
                filename += '.%s' % url_parts[-1]

            urllib.urlretrieve(asset_url, filename)
            use_connector_url = True
        except Exception as e:
            easydb5log(logger,
                       'could not download connector image: %s' % str(e),
                       LOG_WARN)
            return
    else:
        for _file in exp_files:
            if _file['eas_id'] == eas_id:
                filename = '%s/%s' % (exp_files_path, _file['path'])
                break

    if filename is None:
        easydb5log(logger,
                   'no asset file name could be found in insert_picture',
                   LOG_DEBUG)
        return

    try:
        if use_connector_url:
            easydb5log(logger,
                       'load connector image %s' % filename,
                       LOG_DEBUG)
        else:
            easydb5log(logger,
                       'load exported image from local instance %s' % filename,
                       LOG_DEBUG)
        img = Image.open(filename)
    except Exception as e:
        if use_connector_url:
            easydb5log(logger,
                       'could not load connector image: %s' % str(e),
                       LOG_WARN)
        else:
            easydb5log(logger,
                       'could not load exported image from local instance slide: %s' % str(
                           e),
                       LOG_WARN)
        return

    try:
        # get placeholder size in emus
        pw_emu = float(placeholder.width)
        ph_emu = float(placeholder.height)

        iw, ih = img.size

        # convert image size from pixels to emus
        iw_emu = pixels_to_emu(iw)
        ih_emu = pixels_to_emu(ih)

        h_ratio = iw_emu / pw_emu
        w_ratio = ih_emu / ph_emu

        # scale down to fit the longer image side into the shorter placeholder side
        new_x = 0
        new_y = 0
        new_h = 0
        new_w = 0
        if h_ratio >= w_ratio:
            new_h = int(ih_emu / h_ratio)
            new_w = int(iw_emu / h_ratio)
            new_y = (ph_emu - new_h) / 2
        else:
            new_h = int(ih_emu / w_ratio)
            new_w = int(iw_emu / w_ratio)
            new_x = (pw_emu - new_w) / 2

        shapes.add_picture(filename, new_x + placeholder.left,
                           new_y + placeholder.top, height=new_h)

        picture_bottom_line = new_y + placeholder.top + new_h

        # remove the original placeholder since it is not needed
        placeholder._element.getparent().remove(placeholder._element)
    except Exception as e:
        easydb5log(logger,
                   'could not get image resolution / size information, will insert image %s into placeholder' % filename,
                   LOG_WARN)
        placeholder.insert_picture(filename)

    return picture_bottom_line
