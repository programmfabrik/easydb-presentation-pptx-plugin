# encoding: utf-8

import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from PIL import Image
import requests

from fylr_lib_plugin_python3 import util as fylr_util

import hashlib
import os


def __pixels_to_emu(px, dpi=72):
    return float(px * (914400 / dpi))


def __get_standard_format():
    return {
        '1': {
            'size': 26,
            'bold': True
        },
        '2': {
            'size': 19,
            'bold': False
        },
        '3': {
            'size': 16,
            'bold': False
        }
    }


def parse_target_filename(produce_opts):
    return '{0}.pptx'.format(fylr_util.get_json_value(produce_opts, 'presentation.filename', True))


def __create_missing_dirs(f_path):
    base_dir = '/'.join(f_path.split('/')[:-1])
    if not os.path.exists(base_dir):
        os.makedirs(base_dir)


def __parse_slide_layouts(prs, produce_opts, show_standard):
    slide_layouts = {}

    for slide in produce_opts['pptx_form']['template']['slides']:
        if len(show_standard) > 0:
            if 'show_info' in slide and slide['show_info'] == True:
                slide_layouts[slide['type']] = {
                    'layout': prs.slide_layouts[slide['slide_idx']],
                    'info': slide,
                }
        else:
            if 'show_info' not in slide or slide['show_info'] == False:
                slide_layouts[slide['type']] = {
                    'layout': prs.slide_layouts[slide['slide_idx']],
                    'info': slide,
                }

    return slide_layouts


def __insert_info(placeholder, shapes, standard_info, show_standard, standard_format, picture_bottom_line):
    if not isinstance(standard_info, dict):
        return

    top = placeholder.top

    if picture_bottom_line is not None:
        top = min([picture_bottom_line + __pixels_to_emu(10), top])

    text_box = shapes.add_textbox(placeholder.left, top, placeholder.width, placeholder.height)
    text_box.text_frame.word_wrap = True

    first_standard_value = True
    for s in show_standard:
        if not s in standard_format:
            continue
        if not s in standard_info:
            continue

        _s = standard_info[s].strip()
        if len(_s) < 1:
            continue

        if first_standard_value:
            first_standard_value = False
            p = text_box.text_frame.paragraphs[0]
        else:
            p = text_box.text_frame.add_paragraph()

        p.text = _s
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Helvetica'
        p.font.size = Pt(standard_format[s]['size'])
        p.font.bold = standard_format[s]['bold']

    # remove the original placeholder since it is not needed
    placeholder._element.getparent().remove(placeholder._element)


def __insert_text(placeholder, shapes, text):
    if len(text) < 1:
        return

    text_box = shapes.add_textbox(placeholder.left, placeholder.top, placeholder.width, placeholder.height)
    text_box.text_frame.word_wrap = True

    first_line = True
    for s in text.split('\n'):
        s = s.strip()
        if len(s) < 1:
            continue

        if first_line:
            first_line = False
            p = text_box.text_frame.paragraphs[0]
        else:
            p = text_box.text_frame.add_paragraph()

        p.text = s
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.1
        p.font.name = 'Helvetica'
        p.font.size = Pt(26)

    # remove the original placeholder since it is not needed
    placeholder._element.getparent().remove(placeholder._element)


def download_export_file(url, filename):
    resp = requests.get(url)
    if resp.status_code == 200:
        __create_missing_dirs(filename)
        with open(os.path.abspath(filename), 'wb') as outf:
            outf.write(resp.content)
    else:
        raise Exception('could not get file from fylr: status code {0}: {1}'.format(resp.status_code, resp.text))


def __insert_picture(pack_dir, exp_files, placeholder, shapes, eas_id, asset_url):

    if eas_id is None and asset_url is None:
        return None

    picture_bottom_line = None

    filename = None
    use_connector_url = False

    if asset_url is not None:
        try:
            # download the image file, save it in the export asset folder
            m = hashlib.sha1(asset_url)
            filename = '{0}/{1}'.format(pack_dir, str(m.hexdigest()))

            url_parts = asset_url.split('/')
            if len(url_parts) > 1:
                filename += '.{0}'.format(url_parts[-1])

            download_export_file(asset_url, filename)
            use_connector_url = True
        except Exception as e:
            raise Exception('could not download connector image: {0}'.format(str(e)))
    else:
        for _file in exp_files:
            if not 'eas_id' in _file:
                continue
            if _file['eas_id'] != eas_id:
                continue
            filename = os.path.abspath('{0}/{1}'.format(pack_dir, _file['path']))
            break

        if filename is None:
            # no asset for this object
            return picture_bottom_line

    try:
        img = Image.open(filename)
    except Exception as e:
        if use_connector_url:
            raise Exception('could not load connector image: {0}'.format(str(e)))
        else:
            raise Exception('could not load exported image from local instance slide: {0}'.format(str(e)))

    try:
        # get placeholder size in emus
        pw_emu = float(placeholder.width)
        ph_emu = float(placeholder.height)

        iw, ih = img.size

        # convert image size from pixels to emus
        iw_emu = __pixels_to_emu(iw)
        ih_emu = __pixels_to_emu(ih)

        h_ratio = iw_emu / pw_emu
        w_ratio = ih_emu / ph_emu

        # scale down to fit the longer image side into the shorter placeholder side
        new_x = 0
        new_y = 0
        new_h = 0
        new_w = 0
        if h_ratio >= w_ratio:
            new_h = int(ih_emu / h_ratio)
            new_w = int(iw_emu / h_ratio)
            new_y = (ph_emu - new_h) / 2
        else:
            new_h = int(ih_emu / w_ratio)
            new_w = int(iw_emu / w_ratio)
            new_x = (pw_emu - new_w) / 2

        shapes.add_picture(
            filename,
            new_x + placeholder.left,
            new_y + placeholder.top,
            height=new_h
        )

        picture_bottom_line = new_y + placeholder.top + new_h

        # remove the original placeholder since it is not needed
        placeholder._element.getparent().remove(placeholder._element)
    except Exception as e:
        placeholder.insert_picture(filename)

    return picture_bottom_line


def produce_files(produce_opts, pack_dir, export_files, pptx_filename):

    standard_format = __get_standard_format()
    show_standard = fylr_util.get_json_value(produce_opts, 'presentation.settings.show_standard')
    if isinstance(show_standard, str):
        show_standard = list(map(lambda s: s.strip(), show_standard.split()))
    else:
        show_standard = []

    prs = Presentation('{0}/../templates/{1}'.format(
        os.path.abspath(os.path.dirname(__file__)),
        fylr_util.get_json_value(produce_opts, 'pptx_form.template.name', True)))

    slide_layouts = __parse_slide_layouts(prs, produce_opts, show_standard)

    slide_id = -1
    for slide in fylr_util.get_json_value(produce_opts, 'presentation.slides', True):
        slide_id += 1

        stype = slide['type']

        sl = slide_layouts[stype]
        sl_info = sl['info']

        ppt_slide = prs.slides.add_slide(sl['layout'])

        title_key = fylr_util.get_json_value(sl_info, 'title')
        subtitle_key = fylr_util.get_json_value(sl_info, 'subtitle')

        data_title = fylr_util.get_json_value(slide, 'data.title')
        data_info = fylr_util.get_json_value(slide, 'data.info')

        if stype == 'start':
            if not 'data' in slide:
                continue

            if title_key is not None and data_title is not None:
                ppt_slide.placeholders[title_key].text = data_title
            if subtitle_key is not None and data_info is not None:
                ppt_slide.placeholders[subtitle_key].text = data_info

        elif stype == 'bullets':
            if not 'data' in slide:
                continue

            if title_key is not None and data_title is not None:
                ppt_slide.placeholders[title_key].text = data_title

            text_frame = ppt_slide.placeholders[sl_info['bullets']].text_frame
            text_frame.clear()  # remove any existing paragraphs, leaving one empty one

            rows = slide['data']['info'].split('\n')

            p = text_frame.paragraphs[0]
            p.text = rows[0]

            for row in rows[1:]:
                p = text_frame.add_paragraph()
                p.text = row

        elif stype == 'one':
            if not 'center' in slide:
                continue

            if not 'global_object_id' in slide['center']:
                continue

            picture_bottom_line = __insert_picture(
                pack_dir,
                export_files,
                ppt_slide.placeholders[sl_info['picture']],
                ppt_slide.shapes,
                fylr_util.get_json_value(slide, 'center.version_id'),
                fylr_util.get_json_value(slide, 'center.asset_url'))

            if 'text' in sl_info:
                __insert_info(
                    ppt_slide.placeholders[sl_info['text']],
                    ppt_slide.shapes,
                    fylr_util.get_json_value(slide, 'center.standard_info'),
                    show_standard,
                    standard_format,
                    picture_bottom_line)

        elif stype == 'duo':
            picture_bottom_lines = []

            if not 'left' in slide and not 'right' in slide:
                continue

            if 'left' in slide:
                if 'global_object_id' in slide['left'] and 'picture_left' in sl_info:
                    pbl = __insert_picture(
                        pack_dir,
                        export_files,
                        ppt_slide.placeholders[sl_info['picture_left']],
                        ppt_slide.shapes,
                        fylr_util.get_json_value(slide, 'left.version_id'),
                        fylr_util.get_json_value(slide, 'left.asset_url'))
                    if pbl is not None:
                        picture_bottom_lines.append(pbl)

            if 'right' in slide:
                if 'global_object_id' in slide['right'] and 'picture_right' in sl_info:
                    pbl = __insert_picture(
                        pack_dir,
                        export_files,
                        ppt_slide.placeholders[sl_info['picture_right']],
                        ppt_slide.shapes,
                        fylr_util.get_json_value(slide, 'right.version_id'),
                        fylr_util.get_json_value(slide, 'right.asset_url'))
                    if pbl is not None:
                        picture_bottom_lines.append(pbl)

            lowest_picture_bottom_line = None
            if len(picture_bottom_lines) > 0:
                lowest_picture_bottom_line = max(picture_bottom_lines)

            if 'left' in slide:
                if 'global_object_id' in slide['left'] and 'text_left' in sl_info:
                    __insert_info(
                        ppt_slide.placeholders[sl_info['text_left']],
                        ppt_slide.shapes,
                        fylr_util.get_json_value(slide, 'left.standard_info'),
                        show_standard,
                        standard_format,
                        lowest_picture_bottom_line)

            if 'right' in slide:
                if 'global_object_id' in slide['right'] and 'text_right' in sl_info:
                    __insert_info(
                        ppt_slide.placeholders[sl_info['text_right']],
                        ppt_slide.shapes,
                        fylr_util.get_json_value(slide, 'right.standard_info'),
                        show_standard,
                        standard_format,
                        lowest_picture_bottom_line)

        elif stype == 'imageText':
            picture_bottom_lines = []

            if not 'left' in slide and not 'data' in slide:
                continue

            if 'left' in slide:
                if 'global_object_id' in slide['left'] and 'picture_left' in sl_info:
                    picture_bottom_lines.append(__insert_picture(
                        pack_dir,
                        export_files,
                        ppt_slide.placeholders[sl_info['picture_left']],
                        ppt_slide.shapes,
                        fylr_util.get_json_value(slide['left'], 'version_id'),
                        fylr_util.get_json_value(slide['left'], 'asset_url')))

            lowest_picture_bottom_line = None
            picture_bottom_lines = list(filter(None, picture_bottom_lines))
            if len(picture_bottom_lines) > 0:
                lowest_picture_bottom_line = max(picture_bottom_lines)

            if 'left' in slide:
                if 'global_object_id' in slide['left'] and 'text_left' in sl_info:
                    __insert_info(
                        ppt_slide.placeholders[sl_info['text_left']],
                        ppt_slide.shapes,
                        fylr_util.get_json_value(slide, 'left.standard_info'),
                        show_standard,
                        standard_format,
                        lowest_picture_bottom_line)

            text = fylr_util.get_json_value(slide, 'data.text')
            if isinstance(text, str) and 'text_right' in sl_info:
                __insert_text(
                    ppt_slide.placeholders[sl_info['text_right']],
                    ppt_slide.shapes,
                    text)

    __create_missing_dirs(pptx_filename)
    prs.save(pptx_filename)
