
from pptx import Presentation
from pptx.util import Inches
import os

prs1 = Presentation("../src/server/default-black.pptx")
idx = 0

for slide_layout in prs1.slide_layouts:
    print "----", idx
    for shape in slide_layout.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
    idx = idx + 1

exit

prs = Presentation("default.pptx")
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.placeholders[0] # slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

img_path = 'bild_a.jpg'

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(2)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

for i in range(100):
    img_path = 'bild_b.jpg'

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top)

    left = Inches(2)
    height = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)


prs.save('test.pptx')
