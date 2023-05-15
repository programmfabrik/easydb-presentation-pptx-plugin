import os
import hashlib
import urllib
import traceback

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from PIL import Image
from context import get_json_value
import util

LOGGER_NAME = 'pf.plugin.base.presentation-pptx'


def easydb_server_start(easydb_context):
    logger = easydb_context.get_logger(LOGGER_NAME)
    logger.debug('PPTX started')

    easydb_context.register_callback('export_produce', {
        'callback': 'produce_files',
    })


def produce_files(easydb_context, parameters, protocol=None):
    try:
        exp = easydb_context.get_exporter()
        produce_opts = exp.getExport()['export']['produce_options']
        logger = easydb_context.get_logger(LOGGER_NAME)

        if 'pptx' not in produce_opts:
            return

        pack_dir = easydb_context.get_temp_dir()
        pptx_filename = '%s/produce.pptx' % pack_dir
        target_filename = '%s.pptx' % produce_opts['presentation']['filename']

        util.produce_files(
            produce_opts,
            '.',
            export_files,
            pptx_filename)

        exp.addFile(pptx_filename, target_filename)

    except Exception as e:
        traceback.print_exc()
        raise e


# def _produce_files(easydb_context, parameters, protocol):
#     global pack_dir

#     exp = easydb_context.get_exporter()
#     produce_opts = exp.getExport()['export']['produce_options']
#     logger = easydb_context.get_logger(LOGGER_NAME)

#     if 'pptx' not in produce_opts:
#         return

#     if not exp:
#         logger.error('could not get exporter object')
#         return

#     for plugin in easydb_context.get_plugins()['plugins']:
#         if plugin['name'] == 'presentation-pptx':
#             break

#     # produce slides

#     basepath = os.path.abspath(os.path.dirname(__file__))

#     standard_format = {
#         '1': {
#             'size': 26,
#             'bold': True
#         },
#         '2': {
#             'size': 19,
#             'bold': False
#         },
#         '3': {
#             'size': 16,
#             'bold': False
#         }
#     }

#     show_standard = []
#     if 'show_standard' in produce_opts['presentation']['settings']:
#         show_standard = [s.strip() for s in produce_opts['presentation']['settings']['show_standard'].split(' ')]

#     prs = Presentation('%s/%s' % (basepath, produce_opts['pptx_form']['template']['name']))

#     slide_layouts = {}

#     for slide in produce_opts['pptx_form']['template']['slides']:
#         if len(show_standard) > 0:
#             if 'show_info' in slide and slide['show_info'] == True:
#                 slide_layouts[slide['type']] = {
#                     'layout': prs.slide_layouts[slide['slide_idx']],
#                     'info': slide,
#                 }
#         else:
#             if 'show_info' not in slide or slide['show_info'] == False:
#                 slide_layouts[slide['type']] = {
#                     'layout': prs.slide_layouts[slide['slide_idx']],
#                     'info': slide,
#                 }

#     data_by_gid = produce_opts['presentation']['data_by_gid']

#     def insert_info(placeholder, shapes, gid, picture_bottom_line=None):
#         if gid not in data_by_gid:
#             return

#         top = placeholder.top

#         if picture_bottom_line is not None:
#             top = min([picture_bottom_line + pixels_to_emu(10), top])

#         text_box = shapes.add_textbox(placeholder.left, top, placeholder.width, placeholder.height)
#         text_box.text_frame.word_wrap = True

#         first_standard_value = True
#         for s in show_standard:
#             if not s in standard_format:
#                 continue

#             if s in data_by_gid[gid]['standard_info']:
#                 _s = data_by_gid[gid]['standard_info'][s].strip()
#                 if len(_s) < 1:
#                     continue

#                 if first_standard_value:
#                     first_standard_value = False
#                     p = text_box.text_frame.paragraphs[0]
#                 else:
#                     p = text_box.text_frame.add_paragraph()

#                 p.text = _s
#                 p.alignment = PP_ALIGN.CENTER
#                 p.font.name = 'Helvetica'
#                 p.font.size = Pt(standard_format[s]['size'])
#                 p.font.bold = standard_format[s]['bold']

#         # remove the original placeholder since it is not needed
#         placeholder._element.getparent().remove(placeholder._element)

#     def insert_text(placeholder, shapes, text):
#         if len(text) < 1:
#             return

#         text_box = shapes.add_textbox(placeholder.left, placeholder.top, placeholder.width, placeholder.height)
#         text_box.text_frame.word_wrap = True

#         first_line = True
#         for s in text.split('\n'):
#             s = s.strip()
#             if len(s) < 1:
#                 continue

#             if first_line:
#                 first_line = False
#                 p = text_box.text_frame.paragraphs[0]
#             else:
#                 p = text_box.text_frame.add_paragraph()

#             p.text = s
#             p.alignment = PP_ALIGN.LEFT
#             p.line_spacing = 1.1
#             p.font.name = 'Helvetica'
#             p.font.size = Pt(26)

#         # remove the original placeholder since it is not needed
#         placeholder._element.getparent().remove(placeholder._element)

#     def insert_picture(placeholder, shapes, eas_id, asset_url=None):

#         picture_bottom_line = None

#         if eas_id is None and asset_url is None:
#             logger.warn('no asset id or asset url is given for insert_picture')
#             return

#         filename = None
#         use_connector_url = False

#         if asset_url is not None:
#             try:
#                 # download the image file, save it in the export asset folder
#                 m = hashlib.sha1(asset_url)
#                 filename = '%s/%s' % (exp.getFilesPath(), str(m.hexdigest()))

#                 url_parts = asset_url.split('/')
#                 if len(url_parts) > 1:
#                     filename += '.%s' % url_parts[-1]

#                 urllib.urlretrieve(asset_url, filename)
#                 use_connector_url = True
#             except Exception as e:
#                 logger.warn('could not download connector image: %s' % str(e))
#                 return
#         else:
#             for _file in exp.getFiles():
#                 if _file['eas_id'] == eas_id:
#                     filename = '%s/%s' % (exp.getFilesPath(), _file['path'])
#                     break

#         if filename is None:
#             logger.debug('no asset file name could be found in insert_picture')
#             return

#         try:
#             if use_connector_url:
#                 logger.debug('load connector image %s' % filename)
#             else:
#                 logger.debug('load exported image from local instance %s' % filename)
#             img = Image.open(filename)
#         except Exception as e:
#             if use_connector_url:
#                 logger.warn('could not load connector image: %s' % str(e))
#             else:
#                 logger.warn('could not load exported image from local instance slide: %s' % str(e))
#             return

#         try:
#             # get placeholder size in emus
#             pw_emu = float(placeholder.width)
#             ph_emu = float(placeholder.height)

#             iw, ih = img.size

#             # convert image size from pixels to emus
#             iw_emu = pixels_to_emu(iw)
#             ih_emu = pixels_to_emu(ih)

#             h_ratio = iw_emu / pw_emu
#             w_ratio = ih_emu / ph_emu

#             # scale down to fit the longer image side into the shorter placeholder side
#             new_x = 0
#             new_y = 0
#             new_h = 0
#             new_w = 0
#             if h_ratio >= w_ratio:
#                 new_h = int(ih_emu / h_ratio)
#                 new_w = int(iw_emu / h_ratio)
#                 new_y = (ph_emu - new_h) / 2
#             else:
#                 new_h = int(ih_emu / w_ratio)
#                 new_w = int(iw_emu / w_ratio)
#                 new_x = (pw_emu - new_w) / 2

#             shapes.add_picture(filename, new_x + placeholder.left, new_y + placeholder.top, height=new_h)

#             picture_bottom_line = new_y + placeholder.top + new_h

#             # remove the original placeholder since it is not needed
#             placeholder._element.getparent().remove(placeholder._element)
#         except Exception as e:
#             logger.warn('could not get image resolution / size information, will insert image %s into placeholder' % filename)
#             placeholder.insert_picture(filename)

#         return picture_bottom_line

#     slide_id = -1
#     for slide in produce_opts['presentation']['slides']:
#         slide_id += 1

#         stype = slide['type']
#         if stype not in slide_layouts:
#             logger.warn('skipping slide[%d], unknown type: %s' % (slide_id, stype))
#             continue

#         sl = slide_layouts[stype]
#         sl_info = sl['info']

#         logger.debug('adding slide[%d], type: %s | %s | %s' % (slide_id, stype, repr(sl_info), repr(slide)))
#         ppt_slide = prs.slides.add_slide(sl['layout'])

#         if stype == 'start':
#             if not 'data' in slide:
#                 logger.warn('key data missing in slide[%d] in produce_opts' % slide_id)
#                 continue
#             ppt_slide.placeholders[sl_info['title']
#                                    ].text = slide['data']['title']
#             ppt_slide.placeholders[sl_info['subtitle']
#                                    ].text = slide['data']['info']

#         elif stype == 'bullets':
#             if not 'data' in slide:
#                 logger.warn('key data missing in slide[%d] in produce_opts' % slide_id)
#                 continue

#             ppt_slide.placeholders[sl_info['title']].text = slide['data']['title']

#             text_frame = ppt_slide.placeholders[sl_info['bullets']].text_frame
#             text_frame.clear()  # remove any existing paragraphs, leaving one empty one

#             rows = slide['data']['info'].split('\n')

#             p = text_frame.paragraphs[0]
#             p.text = rows[0]

#             for row in rows[1:]:
#                 p = text_frame.add_paragraph()
#                 p.text = row

#         elif stype == 'one':
#             if not 'center' in slide:
#                 logger.warn('key center missing in slide[%d] in produce_opts' % slide_id)
#                 continue

#             if not 'global_object_id' in slide['center']:
#                 logger.warn('key global_object_id missing in slide[%d].center in produce_opts' % slide_id)
#                 continue

#             picture_bottom_line = insert_picture(ppt_slide.placeholders[sl_info['picture']],
#                                                  ppt_slide.shapes,
#                                                  get_json_value(
#                                                      slide['center'], 'asset_id'),
#                                                  get_json_value(slide['center'], 'asset_url'))

#             if 'text' in sl_info:
#                 insert_info(ppt_slide.placeholders[sl_info['text']],
#                             ppt_slide.shapes,
#                             slide['center']['global_object_id'],
#                             picture_bottom_line)

#         elif stype == 'duo':
#             picture_bottom_lines = []

#             if not 'left' in slide and not 'right' in slide:
#                 logger.warn('keys left and right missing in slide[%d] in produce_opts' % slide_id)
#                 continue

#             if 'left' in slide:
#                 if 'global_object_id' in slide['left'] and 'picture_left' in sl_info:
#                     picture_bottom_lines.append(insert_picture(ppt_slide.placeholders[sl_info['picture_left']],
#                                                                ppt_slide.shapes,
#                                                                get_json_value(
#                                                                    slide['left'], 'asset_id'),
#                                                                get_json_value(slide['left'], 'asset_url')))

#             if 'right' in slide:
#                 if 'global_object_id' in slide['right'] and 'picture_right' in sl_info:
#                     picture_bottom_lines.append(insert_picture(ppt_slide.placeholders[sl_info['picture_right']],
#                                                                ppt_slide.shapes,
#                                                                get_json_value(
#                                                                    slide['right'], 'asset_id'),
#                                                                get_json_value(slide['right'], 'asset_url')))

#             lowest_picture_bottom_line = None
#             picture_bottom_lines = list(filter(None, picture_bottom_lines))
#             if len(picture_bottom_lines) > 0:
#                 lowest_picture_bottom_line = max(picture_bottom_lines)

#             if 'left' in slide:
#                 if 'global_object_id' in slide['left'] and 'text_left' in sl_info:
#                     insert_info(ppt_slide.placeholders[sl_info['text_left']],
#                                 ppt_slide.shapes,
#                                 slide['left']['global_object_id'],
#                                 lowest_picture_bottom_line)

#             if 'right' in slide:
#                 if 'global_object_id' in slide['right'] and 'text_right' in sl_info:
#                     insert_info(ppt_slide.placeholders[sl_info['text_right']],
#                                 ppt_slide.shapes,
#                                 slide['right']['global_object_id'],
#                                 lowest_picture_bottom_line)

#         elif stype == 'imageText':
#             picture_bottom_lines = []

#             if not 'left' in slide and not 'data' in slide:
#                 logger.warn('keys left and data missing in slide[%d] in produce_opts' % slide_id)
#                 continue

#             if 'left' in slide:
#                 if 'global_object_id' in slide['left'] and 'picture_left' in sl_info:
#                     picture_bottom_lines.append(insert_picture(ppt_slide.placeholders[sl_info['picture_left']],
#                                                                ppt_slide.shapes,
#                                                                get_json_value(
#                                                                    slide['left'], 'asset_id'),
#                                                                get_json_value(slide['left'], 'asset_url')))

#             lowest_picture_bottom_line = None
#             picture_bottom_lines = list(filter(None, picture_bottom_lines))
#             if len(picture_bottom_lines) > 0:
#                 lowest_picture_bottom_line = max(picture_bottom_lines)

#             if 'left' in slide:
#                 if 'global_object_id' in slide['left'] and 'text_left' in sl_info:
#                     insert_info(ppt_slide.placeholders[sl_info['text_left']],
#                                 ppt_slide.shapes,
#                                 slide['left']['global_object_id'],
#                                 lowest_picture_bottom_line)

#             text = get_json_value(slide, 'data.text')
#             if isinstance(text, str) and 'text_right' in sl_info:
#                 insert_text(ppt_slide.placeholders[sl_info['text_right']],
#                             ppt_slide.shapes,
#                             text)

#         else:
#             logger.warn('unknown type %s in slide[%d] in produce_opts' % (stype, slide_id))

#     pack_dir = easydb_context.get_temp_dir()
#     pptx_filename = '%s/produce.pptx' % pack_dir
#     target_filename = '%s.pptx' % produce_opts['presentation']['filename']

#     prs.save(pptx_filename)
#     exp.addFile(pptx_filename, target_filename)
